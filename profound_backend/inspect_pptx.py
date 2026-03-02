from pptx import Presentation

prs = Presentation('test_output.pptx')
print('slides:', len(prs.slides))
for idx, slide in enumerate(prs.slides, start=1):
    print('slide', idx)
    for shape in slide.shapes:
        if shape.has_text_frame:
            print(' text:', shape.text[:120].replace('\n',' '))
            for para in shape.text_frame.paragraphs:
                try:
                    if para.hyperlink and para.hyperlink.address:
                        print(' link address', para.hyperlink.address)
                except Exception:
                    pass
