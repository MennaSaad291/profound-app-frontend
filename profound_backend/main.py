import os
import bcrypt
import pandas as pd
import io
import json
import re
import urllib.parse
from dotenv import load_dotenv
from urllib.parse import quote_plus
from typing import List, Optional

from fastapi import FastAPI, HTTPException, Depends, UploadFile, File, Form
from fastapi.responses import StreamingResponse
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel, EmailStr
from sqlalchemy import create_engine, Column, Integer, String, ForeignKey, Text
from sqlalchemy.ext.declarative import declarative_base
from sqlalchemy.orm import sessionmaker, Session

from google import genai
from google.genai import types

# PPTX specific imports for styling
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from groq import Groq
import requests

# --- 1. Environment & Database Setup ---
load_dotenv()
db_user = os.getenv("DB_USER")
db_pass = os.getenv("DB_PASSWORD")
db_host = os.getenv("DB_HOST")
db_port = os.getenv("DB_PORT", "5432")
db_name = os.getenv("DB_NAME")

encoded_pass = quote_plus(db_pass) if db_pass else ""
DATABASE_URL = f"postgresql://{db_user}:{encoded_pass}@{db_host}:{db_port}/{db_name}?sslmode=require"

engine = create_engine(
    DATABASE_URL, 
    pool_pre_ping=True, 
    pool_recycle=300,
    connect_args={
        "sslmode": "require", 
        "keepalives": 1, 
        "keepalives_idle": 30, 
        "keepalives_interval": 10, 
        "keepalives_count": 5
    }
)
Base = declarative_base()
SessionLocal = sessionmaker(autocommit=False, autoflush=False, bind=engine)

def get_db():
    db = SessionLocal()
    try: 
        yield db
    finally: 
        db.close()

# --- 2. Database Models ---
class UserDB(Base):
    __tablename__ = "users"
    id = Column(Integer, primary_key=True, index=True)
    full_name = Column(String)
    email = Column(String, unique=True, index=True)
    password_hash = Column(String)
    bio = Column(Text, default="Professor of Computer Science specialized in AI.")
    department = Column(String, default="Information Systems Dept.")

class CourseDB(Base):
    __tablename__ = "courses"
    id = Column(Integer, primary_key=True)
    user_id = Column(Integer, ForeignKey("users.id"))
    code = Column(String)
    name = Column(String)
    semester = Column(String)
    students = Column(Integer, default=0)
    status = Column(String, default="active") 
    schedule = Column(String, default="TBA") 
    room = Column(String, default="TBA") 
    progress = Column(Integer, default=0) 

class StudentDB(Base):
    __tablename__ = "students"
    id = Column(Integer, primary_key=True, index=True)
    student_id = Column(String) 
    name = Column(String)
    department = Column(String)
    course_id = Column(Integer, ForeignKey("courses.id", ondelete="CASCADE"))

class PublicationDB(Base):
    __tablename__ = "publications"
    id = Column(Integer, primary_key=True)
    user_id = Column(Integer, ForeignKey("users.id"))
    title = Column(String)
    journal = Column(String)
    year = Column(Integer)
    citations = Column(Integer, default=0) 

class ProjectDB(Base):
    __tablename__ = "projects"
    id = Column(Integer, primary_key=True)
    user_id = Column(Integer, ForeignKey("users.id"))
    title = Column(String)
    team = Column(String) 
    year = Column(String)
    status = Column(String)

class InterestDB(Base):
    __tablename__ = "interests"
    id = Column(Integer, primary_key=True)
    user_id = Column(Integer, ForeignKey("users.id"))
    name = Column(String)

Base.metadata.create_all(bind=engine)

# --- 3. FastAPI Initialization ---
app = FastAPI(title="Profound Academic API")
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"], 
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# --- 4. Pydantic Schemas ---
class UserCreate(BaseModel): full_name: str; email: EmailStr; password: str
class UserLogin(BaseModel): email: EmailStr; password: str
class UserUpdate(BaseModel): full_name: str; bio: str; department: str
class PublicationCreate(BaseModel): user_id: int; title: str; journal: str; year: int; citations: int = 0
class ProjectCreate(BaseModel): user_id: int; title: str; team: str; year: str; status: str
class InterestCreate(BaseModel): user_id: int; name: str

class CourseResponse(BaseModel):
    id: int; code: str; name: str; semester: str; students: int; status: str; schedule: Optional[str]; room: Optional[str]; progress: Optional[int]
    class Config: from_attributes = True

class LectureGenerationRequest(BaseModel):
    topic: str
    course_level: str
    duration: str
    additional_instructions: str
    include_media: bool
    theme: str

class ExportPPTXRequest(BaseModel):
    slides: list
    theme: str = "Modern Minimalist"

# --- 5. API Endpoints ---
@app.post("/register")
def register_user(user: UserCreate, db: Session = Depends(get_db)):
    if db.query(UserDB).filter(UserDB.email == user.email.lower()).first(): raise HTTPException(status_code=400, detail="Email registered")
    salt = bcrypt.gensalt()
    db.add(UserDB(full_name=user.full_name, email=user.email.lower(), password_hash=bcrypt.hashpw(user.password.encode('utf-8'), salt).decode('utf-8')))
    db.commit()
    return {"message": "Success"}

@app.post("/login")
def login_user(user: UserLogin, db: Session = Depends(get_db)):
    db_user = db.query(UserDB).filter(UserDB.email == user.email.lower()).first()
    if not db_user or not bcrypt.checkpw(user.password.encode('utf-8'), db_user.password_hash.encode('utf-8')): raise HTTPException(status_code=401)
    return {"user": {"id": db_user.id, "name": db_user.full_name}}

@app.get("/profile/{user_id}")
def get_profile(user_id: int, db: Session = Depends(get_db)):
    user = db.query(UserDB).filter(UserDB.id == user_id).first()
    if not user: raise HTTPException(status_code=404)
    pubs = db.query(PublicationDB).filter(PublicationDB.user_id == user_id).all()
    courses = db.query(CourseDB).filter(CourseDB.user_id == user_id).all()
    projects = db.query(ProjectDB).filter(ProjectDB.user_id == user_id).all()
    interests = db.query(InterestDB).filter(InterestDB.user_id == user_id).all()
    return {"id": user.id, "full_name": user.full_name, "bio": user.bio, "department": user.department, "metrics": {"citations": sum(p.citations for p in pubs), "students": sum(c.students for c in courses), "papers": len(pubs), "projects": len(projects)}, "publications": pubs, "courses": courses, "projects": projects, "interests": [i.name for i in interests]}

@app.post("/profile/update/{user_id}")
def update_profile(user_id: int, update_data: UserUpdate, db: Session = Depends(get_db)):
    db_user = db.query(UserDB).filter(UserDB.id == user_id).first()
    db_user.full_name, db_user.bio, db_user.department = update_data.full_name, update_data.bio, update_data.department
    db.commit()
    return {"message": "Success"}

@app.post("/courses-with-students")
async def create_course_with_excel(user_id: int = Form(...), code: str = Form(...), name: str = Form(...), semester: str = Form("TBA"), schedule: str = Form("TBA"), room: str = Form("TBA"), file: Optional[UploadFile] = File(None), db: Session = Depends(get_db)):
    student_count, df = 0, None
    if file and file.filename:
        df = pd.read_excel(io.BytesIO(await file.read()))
        student_count = len(df)
    new_course = CourseDB(user_id=user_id, code=code, name=name, semester=semester, students=student_count, status="active", schedule=schedule, room=room, progress=0)
    db.add(new_course)
    db.flush() 
    if df is not None:
        for _, row in df.iterrows(): db.add(StudentDB(student_id=str(row['id']), name=row['name'], department=row.get('department', 'N/A'), course_id=new_course.id))
    db.commit()
    return {"message": "Success"}

@app.get("/professors/{user_id}/courses", response_model=List[CourseResponse])
def get_courses(user_id: int, db: Session = Depends(get_db)): return db.query(CourseDB).filter(CourseDB.user_id == user_id).all()

@app.delete("/courses/{course_id}")
def delete_course(course_id: int, db: Session = Depends(get_db)):
    db_course = db.query(CourseDB).filter(CourseDB.id == course_id).first()
    db.delete(db_course)
    db.commit()
    return {"message": "Deleted"}


import os
import io
import json
import re
import requests
import urllib.parse
from fastapi import FastAPI, HTTPException
from fastapi.responses import StreamingResponse
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from groq import Groq
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE

app = FastAPI()
app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_methods=["*"], allow_headers=["*"])

client = Groq(api_key=os.getenv("GROQ_API_KEY"))

THEME_COLORS = {
    "Modern Minimalist": {"bg": "FFFFFF", "text": "0F172A", "accent": "9333EA", "accent2": "EC4899"},
    "Dark Mode Tech": {"bg": "1E293B", "text": "FFFFFF", "accent": "38BDF8", "accent2": "06B6D4"},
    "Classic Academic": {"bg": "FDFBF7", "text": "0F172A", "accent": "8B1A1A", "accent2": "D4844F"},
    "Vibrant Creative": {"bg": "FFF7ED", "text": "0F172A", "accent": "F97316", "accent2": "EC4899"},
}

class LectureRequest(BaseModel):
    topic: str
    course_level: str
    pages_count: int
    additional_instructions: str
    include_media: bool
    theme: str

def clean_markdown(text: str) -> str:
    return text.replace('**', '').replace('*', '').replace('##', '').replace('#', '').strip()

@app.post("/api/generate-lecture")
async def generate_lecture(data: LectureRequest):
    prompt = f"""
    Generate a university lecture JSON for topic: {data.topic}.
    Slide Count: {data.pages_count}.
    Instructions: {data.additional_instructions}
    
    Format:
    {{ "slides": [ {{ "title": "...", "content": ["Paragraph explanation 1", "Paragraph 2"], "speaker_notes": "..." }} ] }}
    """
    
    completion = client.chat.completions.create(
        model="llama-3.1-8b-instant",
        messages=[{"role": "system", "content": "You are a JSON-only academic generator."}, {"role": "user", "content": prompt}],
        temperature=0.7,
    )
    
    raw_content = completion.choices[0].message.content
    # Clean up the AI output to ensure valid JSON
    json_str = re.search(r'\{.*\}', raw_content, re.DOTALL).group(0)
    
    try:
        return json.loads(json_str)
    except Exception as e:
        raise HTTPException(status_code=500, detail="AI response was not valid JSON.")

@app.post("/api/export-pptx")
async def export_pptx(data: dict):
    theme = data.get('theme', 'Modern Minimalist')
    colors = THEME_COLORS.get(theme, THEME_COLORS['Modern Minimalist'])
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5) # 16:9

    for slide_data in data.get('slides', []):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(colors['bg'])

        # Add decorative "ProFound" style background circles
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(-1), Inches(4), Inches(4))
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor.from_string(colors['accent'])
        circle.line.fill.background()

        # Title Styling
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
        tf = title_box.text_frame
        tf.paragraphs[0].text = clean_markdown(slide_data.get("title", ""))
        tf.paragraphs[0].font.size = Pt(40)
        tf.paragraphs[0].font.color.rgb = RGBColor.from_string(colors['accent'])

        # Body Content Styling
        body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(4.5))
        btf = body_box.text_frame
        btf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        for point in slide_data.get("content", []):
            p = btf.add_paragraph()
            p.text = f"• {clean_markdown(point)}"
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor.from_string(colors['text'])

    stream = io.BytesIO()
    prs.save(stream)
    stream.seek(0)
    return StreamingResponse(stream, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")